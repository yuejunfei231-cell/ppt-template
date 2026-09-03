# -*- coding: utf-8 -*-
"""deckkit.to_pptx —— 把 Canvas 渲染成可编辑的 PowerPoint (.pptx)。

生成的都是原生形状与文本框（矢量），可在 PowerPoint / WPS / Keynote 里直接改字改色。
"""
from __future__ import annotations

from typing import Any, Dict, List, Sequence

from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from . import theme as T
from .layout import Canvas, PT2EMU, Para, Run

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT,
          "j": PP_ALIGN.JUSTIFY}
_ANCHOR = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}
_DASH = {"dash": "dash", "dot": "sysDot", "dashdot": "dashDot", "lgdash": "lgDash"}


def E(v: float) -> int:
    return int(round(v * PT2EMU))


def _rgb(hexstr: str) -> RGBColor:
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _sub(parent, tag, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def _paint(shape, fill=None, line=None, lw=0.75, alpha=None, dash=None,
           grad=None, cap=None):
    sp = shape
    sp.shadow.inherit = False
    # 填充
    if grad:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(grad[0])
        sF = sp._element.spPr.find(qn("a:solidFill"))
        g = _sub(sp._element.spPr, "a:gradFill", rotWithShape="1")
        gsLst = _sub(g, "a:gsLst")
        for pos, col in ((0, grad[0]), (100000, grad[1])):
            gs = _sub(gsLst, "a:gs", pos=pos)
            _sub(_sub(gs, "a:srgbClr", val=col), "a:alpha", val=100000)
        _sub(g, "a:lin", ang=int(grad[2] * 60000), scaled="1")
        sF.getparent().remove(sF)
        sp._element.spPr.remove(g)
        spPr = sp._element.spPr
        geom = spPr.find(qn("a:prstGeom"))
        if geom is None:
            geom = spPr.find(qn("a:custGeom"))
        geom.addnext(g)
    elif fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        if alpha is not None and alpha < 1.0:
            clr = sp._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
            _sub(clr, "a:alpha", val=int(alpha * 100000))
    else:
        sp.fill.background()
    # 描边
    if line:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(lw)
        ln = sp._element.spPr.find(qn("a:ln"))
        if alpha is not None and alpha < 1.0:
            sc = ln.find(qn("a:solidFill"))
            if sc is not None:
                clr = sc.find(qn("a:srgbClr"))
                if clr is not None:
                    _sub(clr, "a:alpha", val=int(alpha * 100000))
        if dash:
            for e in ln.findall(qn("a:prstDash")):
                ln.remove(e)
            _sub(ln, "a:prstDash", val=_DASH.get(dash, "dash"))
        if cap:
            ln.set("cap", cap)
    else:
        sp.line.fill.background()


def _shadow_on(shape, blur=5.0, dist=2.0, dir_deg=90.0, color="0A1E34", alpha=0.24):
    spPr = shape._element.spPr
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = _sub(spPr, "a:effectLst")
    sh = _sub(eff, "a:outerShdw", blurRad=E(blur), dist=E(dist),
              dir=int(dir_deg * 60000), rotWithShape="0")
    c = _sub(sh, "a:srgbClr", val=color)
    _sub(c, "a:alpha", val=int(alpha * 100000))


def _freeform(slide, pts, close=True):
    shapes = slide.shapes
    fb = shapes.build_freeform(E(pts[0][0]), E(pts[0][1]), scale=1.0)
    fb.add_line_segments([(E(x), E(y)) for x, y in pts[1:]], close=close)
    return fb.convert_to_shape()


def _run_xml(run_el, style: Run):
    run_el.text = style.text
    f = run_el.font
    f.size = Pt(style.size)
    f.bold = style.bold
    f.italic = style.italic
    f.color.rgb = _rgb(style.color)
    f.name = style.font_latin
    rPr = run_el._r.get_or_add_rPr()
    for tag, face in (("a:latin", style.font_latin), ("a:ea", style.font_ea),
                      ("a:cs", style.font_latin)):
        e = rPr.find(qn(tag))
        if e is None:
            e = _sub(rPr, tag)
        e.set("typeface", face)
    if style.baseline:
        rPr.set("baseline", str(int(style.baseline * 100000)))


def _emit_paras(tf, paras: Sequence[Para], valign="t"):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _ANCHOR.get(valign, MSO_ANCHOR.TOP)
    for i, p in enumerate(paras):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = _ALIGN.get(p.align, PP_ALIGN.LEFT)
        par.line_spacing = p.spacing
        if p.space_before:
            par.space_before = Pt(p.space_before)
        if p.space_after:
            par.space_after = Pt(p.space_after)
        pPr = par._p.get_or_add_pPr()
        if p.bullet:
            marL = p.indent + p.bullet_gap
            pPr.set("marL", str(E(marL)))
            pPr.set("indent", str(E(-p.bullet_gap)))
            bc = _sub(pPr, "a:buClr")
            _sub(bc, "a:srgbClr", val=(p.bullet_color or T.GRAY))
            _sub(pPr, "a:buFont", typeface="Arial", pitchFamily="34", charset="0")
            _sub(pPr, "a:buChar", char=p.bullet)
        elif p.indent:
            pPr.set("marL", str(E(p.indent)))
            pPr.set("indent", "0")
        for r in p.runs:
            _run_xml(par.add_run(), r)


def render(cv: Canvas, prs: Presentation):
    blank = None
    for lay in prs.slide_layouts:
        if lay.name.lower() in ("blank", "空白"):
            blank = lay
            break
    if blank is None:
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    # 背景
    if cv.bg:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(cv.bg)
    shapes = slide.shapes
    for it in cv.items:
        k = it["kind"]
        if k == "rect":
            r = it.get("radius", 0)
            shp = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE,
                E(it["x"]), E(it["y"]), E(it["w"]), E(it["h"]))
            if r:
                try:
                    shp.adjustments[0] = min(0.5, r / max(1e-6, min(it["w"], it["h"])))
                except Exception:
                    pass
            _paint(shp, it.get("fill"), it.get("line"), it.get("lw", 0.75),
                   it.get("alpha"), it.get("dash"), it.get("grad"))
            if it.get("shadow"):
                _shadow_on(shp)
        elif k == "ellipse":
            shp = shapes.add_shape(MSO_SHAPE.OVAL, E(it["x"]), E(it["y"]),
                                   E(it["w"]), E(it["h"]))
            _paint(shp, it.get("fill"), it.get("line"), it.get("lw", 0.75),
                   it.get("alpha"), it.get("dash"))
            if it.get("shadow"):
                _shadow_on(shp)
        elif k in ("poly", "path", "line"):
            pts = it["pts"]
            if len(pts) < 2:
                continue
            shp = _freeform(slide, pts, close=it.get("close", k == "poly"))
            _paint(shp, it.get("fill"), it.get("line"), it.get("lw", 0.75),
                   it.get("alpha"), it.get("dash"), cap=it.get("cap"))
            if it.get("shadow"):
                _shadow_on(shp)
        elif k == "image":
            pic = shapes.add_picture(it["path"], E(it["x"]), E(it["y"]),
                                     E(it["w"]), E(it["h"]))
            pic.shadow.inherit = False
            if it.get("shadow"):
                _shadow_on(pic)
        elif k == "text":
            h = max(1.0, it.get("h") or 20)
            tb = shapes.add_textbox(E(it["x"]), E(it["y"]), E(it["w"]), E(h))
            _emit_paras(tb.text_frame, it["paras"], it.get("valign", "t"))
    return slide


def new_deck(title: str, author: str = "", subject: str = "", comments: str = "") -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(E(T.SLIDE_W))
    prs.slide_height = Emu(E(T.SLIDE_H))
    cp = prs.core_properties
    cp.title = title
    cp.author = author
    cp.subject = subject
    cp.comments = comments
    return prs
